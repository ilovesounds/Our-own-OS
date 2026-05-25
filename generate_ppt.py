import sys
import shutil
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    # Source paths for generated images
    img_core_src = r"C:\Users\soham\.gemini\antigravity\brain\39af0efe-7300-4982-9a4c-e8d9f6b7db3a\kratos_os_core_network_1779710437359.png"
    img_broken_src = r"C:\Users\soham\.gemini\antigravity\brain\39af0efe-7300-4982-9a4c-e8d9f6b7db3a\broken_code_glitch_1779710459784.png"
    img_clock_src = r"C:\Users\soham\.gemini\antigravity\brain\39af0efe-7300-4982-9a4c-e8d9f6b7db3a\neon_timeline_clock_1779710481000.png"

    # Local workspace destination paths
    img_core = "core_network.png"
    img_broken = "broken_code.png"
    img_clock = "timeline_clock.png"

    # Copy files locally to root workspace
    try:
        shutil.copy(img_core_src, img_core)
        shutil.copy(img_broken_src, img_broken)
        shutil.copy(img_clock_src, img_clock)
        print("Successfully copied slide assets to workspace root.")
    except Exception as e:
        print(f"Warning: Could not copy assets: {e}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors
    c_bg = RGBColor(10, 10, 12)          # Deep Cyber Black
    c_card = RGBColor(20, 20, 24)        # Dark Card Gray
    c_border = RGBColor(39, 39, 42)      # Zinc 800 Border
    c_lime = RGBColor(163, 230, 53)      # Lime 400 Accent
    c_cyan = RGBColor(6, 182, 212)      # Cyan 400 Accent
    c_red = RGBColor(239, 68, 68)        # Red 500 Chaos
    c_white = RGBColor(255, 255, 255)    # Plain White
    c_gray = RGBColor(161, 161, 170)     # Zinc 400 Text
    c_dark_gray = RGBColor(82, 82, 91)   # Zinc 600 Info

    def set_bg(slide):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = c_bg

    def draw_card(slide, left, top, width, height, bg_color=c_card, border_color=c_border):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
        return shape

    def add_textbox(slide, left, top, width, height, text="", font_name="Calibri", font_size=12, font_color=c_white, bold=False, align=PP_ALIGN.LEFT, word_wrap=True):
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.word_wrap = word_wrap
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.05)
        
        p = tf.paragraphs[0]
        p.alignment = align
        if text:
            p.text = text
            p.font.name = font_name
            p.font.size = Pt(font_size)
            p.font.color.rgb = font_color
            p.font.bold = bold
        return tx

    def add_bullet(tf, text, font_name="Calibri", font_size=12, font_color=c_white, bold=False, level=0):
        p = tf.add_paragraph()
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.level = level

    def add_slide_header(slide, title_text, category="KRATOS OS ARCHITECTURE"):
        # Dark header bar
        draw_card(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.1), c_card, c_border)
        # Category label
        add_textbox(slide, Inches(0.6), Inches(0.18), Inches(8), Inches(0.3), category, "Courier New", 9, c_lime, True)
        # Slide Title
        add_textbox(slide, Inches(0.6), Inches(0.42), Inches(8), Inches(0.55), title_text, "Consolas", 22, c_white, True)
        # Logo badge in top right
        draw_card(slide, Inches(10.8), Inches(0.35), Inches(1.933), Inches(0.4), c_bg, c_border)
        add_textbox(slide, Inches(10.8), Inches(0.42), Inches(1.933), Inches(0.3), "KRATOS OS V1.0", "Courier New", 9, c_lime, True, PP_ALIGN.CENTER)

    def add_slide_footer(slide, current_index):
        add_textbox(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3), "KRATOS OS SANDBOX SYSTEMS", "Courier New", 9, c_dark_gray)
        add_textbox(slide, Inches(8.5), Inches(7.05), Inches(3), Inches(0.3), "CONFIDENTIAL // DEVONLY", "Courier New", 9, c_dark_gray, False, PP_ALIGN.RIGHT)
        add_textbox(slide, Inches(11.8), Inches(7.05), Inches(1), Inches(0.3), f"SLIDE {current_index:02d} / 12", "Courier New", 9, c_lime, True, PP_ALIGN.RIGHT)

    def try_add_picture(slide, path, left, top, width, height):
        if os.path.exists(path):
            try:
                slide.shapes.add_picture(path, left, top, width, height)
            except Exception as e:
                print(f"Error adding image {path}: {e}")
        else:
            print(f"Image not found: {path}")

    # ----------------------------------------------------
    # SLIDE 1: TITLE SLIDE (WITH COVER IMAGE)
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1)
    # Background card decoration
    draw_card(s1, Inches(0.5), Inches(0.5), Inches(12.333), Inches(6.5), c_card, c_border)
    
    # Left Text Block
    add_textbox(s1, Inches(1.0), Inches(1.8), Inches(6.0), Inches(1.2), "KRATOS OS", "Consolas", 52, c_lime, True)
    add_textbox(s1, Inches(1.0), Inches(3.0), Inches(6.0), Inches(1.0), "The Self-Healing Multi-Agentic\nSandbox Environment", "Calibri", 20, c_white, False)
    add_textbox(s1, Inches(1.0), Inches(4.3), Inches(6.0), Inches(0.4), "PROMPT | DEVCORE | CHAOS | CHRONOS", "Courier New", 11, c_cyan, True)
    add_textbox(s1, Inches(1.0), Inches(5.6), Inches(6.0), Inches(0.4), "LANGGRAPH ENGINE  •  HYDRADB TIMELINE LEDGER", "Courier New", 9, c_dark_gray)

    # Right Image Block (Core network)
    img_l, img_t, img_w, img_h = Inches(7.4), Inches(1.3), Inches(4.8), Inches(4.8)
    draw_card(s1, img_l, img_t, img_w, img_h, c_bg, c_border)
    try_add_picture(s1, img_core, img_l + Inches(0.05), img_t + Inches(0.05), img_w - Inches(0.1), img_h - Inches(0.1))

    add_slide_footer(s1, 1)

    # ----------------------------------------------------
    # SLIDE 2: THE CHALLENGE (WITH BROKEN CODE IMAGE)
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2)
    add_slide_header(s2, "The Problem: Failure in Agentic Systems", "THE OPERATIONAL HAZARDS OF LLM AGENTS")
    
    # Left Column - Text explanation
    draw_card(s2, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), c_card, c_border)
    tx_box = add_textbox(s2, Inches(0.85), Inches(1.7), Inches(5.3), Inches(4.8))
    tf = tx_box.text_frame
    tf.paragraphs[0].text = "Operational Vulnerabilities in AI Runtimes:"
    tf.paragraphs[0].font.name = "Consolas"
    tf.paragraphs[0].font.size = Pt(15)
    tf.paragraphs[0].font.color.rgb = c_white
    tf.paragraphs[0].font.bold = True
    
    add_bullet(tf, "Silent Logic Failures: LLMs frequently generate broken parameters or route incorrectly without firing stderr exits.", "Calibri", 12, c_gray)
    add_bullet(tf, "Missing Agentic Telemetry: Multi-agent message loops and state dependencies are obscured in standard log streams.", "Calibri", 12, c_gray)
    add_bullet(tf, "Irreversible State Mutations: File writes and workspace alterations cannot be recovered without full git rollbacks.", "Calibri", 12, c_gray)
    add_bullet(tf, "Lack of Automated Recovery: A compilation crash completely interrupts pipeline tasks, requiring human repair.", "Calibri", 12, c_gray)

    # Right Column - Image (Top) & Code Mockup (Bottom)
    draw_card(s2, Inches(6.8), Inches(1.5), Inches(5.9), Inches(2.4), c_bg, c_border)
    try_add_picture(s2, img_broken, Inches(6.85), Inches(1.55), Inches(5.8), Inches(2.3))

    draw_card(s2, Inches(6.8), Inches(4.1), Inches(5.9), Inches(2.6), RGBColor(15, 10, 10), c_red)
    add_textbox(s2, Inches(7.1), Inches(4.2), Inches(5.3), Inches(0.35), "// CRITICAL THREAD EXCEPTION TRACE", "Courier New", 10, c_red, True)
    code_box = add_textbox(s2, Inches(7.1), Inches(4.6), Inches(5.3), Inches(2.0))
    code_tf = code_box.text_frame
    code_tf.paragraphs[0].text = "File \"db_connection.py\", line 12, in get_db_client"
    code_tf.paragraphs[0].font.name = "Courier New"
    code_tf.paragraphs[0].font.size = Pt(9)
    code_tf.paragraphs[0].font.color.rgb = c_gray
    add_bullet(code_tf, "  raise ConnectionError(\"HydraDB socket timeout\")", "Courier New", 9, c_gray)
    add_bullet(code_tf, "ConnectionError: HydraDB socket timeout (code 504)", "Courier New", 9, c_red, True)
    add_bullet(code_tf, "[FATAL] Thread loop terminated. Workspace locked.", "Courier New", 9, c_dark_gray)

    add_slide_footer(s2, 2)

    # ----------------------------------------------------
    # SLIDE 3: THE PARADIGM
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3)
    add_slide_header(s3, "The Kratos OS Paradigm: Core Orchestration", "CLOSED-LOOP AGENT SYSTEMS")
    
    card_w = Inches(3.77)
    card_h = Inches(5.2)
    spacing = Inches(0.4)
    start_x = Inches(0.6)
    
    # Card 1
    draw_card(s3, start_x, Inches(1.5), card_w, card_h, c_card, c_border)
    add_textbox(s3, start_x + Inches(0.2), Inches(1.7), card_w - Inches(0.4), Inches(0.4), "01 / INTENT PARSING", "Courier New", 12, c_lime, True)
    c1_txt = add_textbox(s3, start_x + Inches(0.2), Inches(2.2), card_w - Inches(0.4), Inches(4.2))
    c1_tf = c1_txt.text_frame
    c1_tf.paragraphs[0].text = "Structured Goal Building:"
    c1_tf.paragraphs[0].font.name = "Consolas"
    c1_tf.paragraphs[0].font.size = Pt(13)
    c1_tf.paragraphs[0].font.color.rgb = c_white
    c1_tf.paragraphs[0].font.bold = True
    add_bullet(c1_tf, "Translates messy conversational language into executable JSON steps.", "Calibri", 11, c_gray)
    add_bullet(c1_tf, "Evaluates file dependencies and locks correct workspace nodes before operations.", "Calibri", 11, c_gray)
    add_bullet(c1_tf, "Establishes automated unit test boundaries to check code correctness.", "Calibri", 11, c_gray)

    # Card 2
    draw_card(s3, start_x + card_w + spacing, Inches(1.5), card_w, card_h, c_card, c_border)
    add_textbox(s3, start_x + card_w + spacing + Inches(0.2), Inches(1.7), card_w - Inches(0.4), Inches(0.4), "02 / CLOSED-LOOP CORE", "Courier New", 12, c_cyan, True)
    c2_txt = add_textbox(s3, start_x + card_w + spacing + Inches(0.2), Inches(2.2), card_w - Inches(0.4), Inches(4.2))
    c2_tf = c2_txt.text_frame
    c2_tf.paragraphs[0].text = "Self-Healing Runtimes:"
    c2_tf.paragraphs[0].font.name = "Consolas"
    c2_tf.paragraphs[0].font.size = Pt(13)
    c2_tf.paragraphs[0].font.color.rgb = c_white
    c2_tf.paragraphs[0].font.bold = True
    add_bullet(c2_tf, "Integrates Planner, Coder, and Tester in a continuous execution cycle.", "Calibri", 11, c_gray)
    add_bullet(c2_tf, "If verification fails, the Corrector Core rewrites the target module immediately.", "Calibri", 11, c_gray)
    add_bullet(c2_tf, "Prevents syntax error propagation to active user endpoints.", "Calibri", 11, c_gray)

    # Card 3
    draw_card(s3, start_x + (card_w + spacing)*2, Inches(1.5), card_w, card_h, c_card, c_border)
    add_textbox(s3, start_x + (card_w + spacing)*2 + Inches(0.2), Inches(1.7), card_w - Inches(0.4), Inches(0.4), "03 / LEDGER ROLLBACK", "Courier New", 12, c_white, True)
    c3_txt = add_textbox(s3, start_x + (card_w + spacing)*2 + Inches(0.2), Inches(2.2), card_w - Inches(0.4), Inches(4.2))
    c3_tf = c3_txt.text_frame
    c3_tf.paragraphs[0].text = "Immutable Snapshots:"
    c3_tf.paragraphs[0].font.name = "Consolas"
    c3_tf.paragraphs[0].font.size = Pt(13)
    c3_tf.paragraphs[0].font.color.rgb = c_white
    c3_tf.paragraphs[0].font.bold = True
    add_bullet(c3_tf, "Saves full workspace state checks at every agent state transition.", "Calibri", 11, c_gray)
    add_bullet(c3_tf, "Enables absolute time-travel recovery of code buffers and checklists.", "Calibri", 11, c_gray)
    add_bullet(c3_tf, "Uses HydraDB client integration to secure and verify system journals.", "Calibri", 11, c_gray)

    add_slide_footer(s3, 3)

    # ----------------------------------------------------
    # SLIDE 4: THE 6-CORE ENGINE (CORES 00 & 01)
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4)
    add_slide_header(s4, "6-Core Processor: Planning & Coding", "CORES 00 AND 01")
    
    # Left Card: Core 00
    draw_card(s4, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), c_card, c_border)
    add_textbox(s4, Inches(0.9), Inches(1.7), Inches(5.2), Inches(0.4), "CORE 00 // PLANNER", "Courier New", 13, c_lime, True)
    c0_txt = add_textbox(s4, Inches(0.9), Inches(2.2), Inches(5.2), Inches(4.2))
    c0_tf = c0_txt.text_frame
    c0_tf.paragraphs[0].text = "Checklist Generation and Dependency Mapping"
    c0_tf.paragraphs[0].font.name = "Consolas"
    c0_tf.paragraphs[0].font.size = Pt(14)
    c0_tf.paragraphs[0].font.color.rgb = c_white
    c0_tf.paragraphs[0].font.bold = True
    add_bullet(c0_tf, "Role: Translates user intent into step-by-step checklists.", "Calibri", 12, c_gray)
    add_bullet(c0_tf, "Parsing: Breaks down file structures, routing targets, and specific verification requirements.", "Calibri", 12, c_gray)
    add_bullet(c0_tf, "Integrations: Formulates validation conditions that Core 02 must test against.", "Calibri", 12, c_gray)
    add_bullet(c0_tf, "Output: JSON action sequence outlining file modifications.", "Calibri", 12, c_gray)

    # Right Card: Core 01
    draw_card(s4, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.2), c_card, c_border)
    add_textbox(s4, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.4), "CORE 01 // CODER", "Courier New", 13, c_lime, True)
    c1_txt2 = add_textbox(s4, Inches(7.1), Inches(2.2), Inches(5.3), Inches(4.2))
    c1_tf2 = c1_txt2.text_frame
    c1_tf2.paragraphs[0].text = "Source Code Construction and File Buffering"
    c1_tf2.paragraphs[0].font.name = "Consolas"
    c1_tf2.paragraphs[0].font.size = Pt(14)
    c1_tf2.paragraphs[0].font.color.rgb = c_white
    c1_tf2.paragraphs[0].font.bold = True
    add_bullet(c1_tf2, "Role: Receives steps from the Planner and writes source code.", "Calibri", 12, c_gray)
    add_bullet(c1_tf2, "Tasks: Autonomously scaffolds layouts, styles components, and implements backend FastAPI endpoint logics.", "Calibri", 12, c_gray)
    add_bullet(c1_tf2, "Write Cycle: Directly modifies files within the project buffer, preserving unchanged methods.", "Calibri", 12, c_gray)
    add_bullet(c1_tf2, "Target: Generates both server-side logic and frontend interfaces.", "Calibri", 12, c_gray)

    add_slide_footer(s4, 4)

    # ----------------------------------------------------
    # SLIDE 5: THE 6-CORE ENGINE (CORES 02 - 05)
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5)
    add_slide_header(s5, "6-Core Processor: Runtime, Self-Healing & Ledger", "CORES 02 TO 05")
    
    gw = Inches(5.8)
    gh = Inches(2.4)
    
    # Core 02 - Top Left
    draw_card(s5, Inches(0.6), Inches(1.5), gw, gh, c_card, c_border)
    add_textbox(s5, Inches(0.8), Inches(1.6), gw - Inches(0.4), Inches(0.3), "CORE 02 // TESTER", "Courier New", 11, c_cyan, True)
    tx = add_textbox(s5, Inches(0.8), Inches(1.95), gw - Inches(0.4), gh - Inches(0.55))
    tf = tx.text_frame
    tf.paragraphs[0].text = "Runs automated checks, verifies compilation status, and asserts outputs. Translates validation errors to clear corrector guidelines."
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = c_gray

    # Core 03 - Top Right
    draw_card(s5, Inches(6.8), Inches(1.5), gw, gh, c_card, c_border)
    add_textbox(s5, Inches(7.0), Inches(1.6), gw - Inches(0.4), Inches(0.3), "CORE 03 // CORRECTOR", "Courier New", 11, c_cyan, True)
    tx = add_textbox(s5, Inches(7.0), Inches(1.95), gw - Inches(0.4), gh - Inches(0.55))
    tf = tx.text_frame
    tf.paragraphs[0].text = "Captures stdout/stderr traces from Core 02 failures. Parses traceback line exceptions, edits the broken segment, and schedules a compiler recheck."
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = c_gray

    # Core 04 - Bottom Left
    draw_card(s5, Inches(0.6), Inches(4.3), gw, gh, c_card, c_border)
    add_textbox(s5, Inches(0.8), Inches(4.4), gw - Inches(0.4), Inches(0.3), "CORE 04 // DEPLOYER", "Courier New", 11, c_cyan, True)
    tx = add_textbox(s5, Inches(0.8), Inches(4.75), gw - Inches(0.4), gh - Inches(0.55))
    tf = tx.text_frame
    tf.paragraphs[0].text = "Takes successful code buffers, manages system assets, executes production builds, restarts the web kernel, and updates active UI components."
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = c_gray

    # Core 05 - Bottom Right
    draw_card(s5, Inches(6.8), Inches(4.3), gw, gh, c_card, c_border)
    add_textbox(s5, Inches(7.0), Inches(4.4), gw - Inches(0.4), Inches(0.3), "CORE 05 // LEDGER", "Courier New", 11, c_cyan, True)
    tx = add_textbox(s5, Inches(7.0), Inches(4.75), gw - Inches(0.4), gh - Inches(0.55))
    tf = tx.text_frame
    tf.paragraphs[0].text = "Persists execution history. Records planner status, snapshot hashes, file buffers, and timestamps in an immutable local timeline ledger."
    tf.paragraphs[0].font.name = "Calibri"
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = c_gray

    add_slide_footer(s5, 5)

    # ----------------------------------------------------
    # SLIDE 6: MODE 01: PROMPT MODE
    # ----------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6)
    add_slide_header(s6, "Mode 01: Prompt Mode (Intent Compiler)", "SYSTEM MODES OVERVIEW")
    
    # Left Column
    draw_card(s6, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), c_card, c_border)
    tx = add_textbox(s6, Inches(0.9), Inches(1.7), Inches(5.2), Inches(4.8))
    tf = tx.text_frame
    tf.paragraphs[0].text = "Natural Language Action Compiler"
    tf.paragraphs[0].font.name = "Consolas"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = c_white
    tf.paragraphs[0].font.bold = True
    add_bullet(tf, "Input Channel: Developers interact with a sleek, conversational dashboard.", "Calibri", 12, c_gray)
    add_bullet(tf, "LLM Orchestrator: Leverages Llama 3.3 and Claude 3.5 Sonnet to map raw goals to structured schemas.", "Calibri", 12, c_gray)
    add_bullet(tf, "Schema Output: Generates target file paths, parameter limits, and core routing logic.", "Calibri", 12, c_gray)
    add_bullet(tf, "Task Queuing: Automatically sets core co-processor registers to PLANNING or WRITING to begin work.", "Calibri", 12, c_gray)

    # Right Column - JSON Mock
    draw_card(s6, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.2), RGBColor(16, 18, 20), c_border)
    add_textbox(s6, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.4), "// COMPILED ACTION SCHEMA (PROMPT_MODE)", "Courier New", 11, c_cyan, True)
    
    json_box = add_textbox(s6, Inches(7.1), Inches(2.2), Inches(5.3), Inches(4.2))
    json_tf = json_box.text_frame
    json_tf.word_wrap = True
    json_tf.paragraphs[0].text = "{"
    json_tf.paragraphs[0].font.name = "Courier New"
    json_tf.paragraphs[0].font.size = Pt(10.5)
    json_tf.paragraphs[0].font.color.rgb = c_gray
    add_bullet(json_tf, '  "intent": "Scaffold a logistics status router"', "Courier New", 10.5, c_cyan)
    add_bullet(json_tf, '  "target_file": "logistics_tracker.py",', "Courier New", 10.5, c_gray)
    add_bullet(json_tf, '  "action_type": "GENERATE_CODE",', "Courier New", 10.5, c_gray)
    add_bullet(json_tf, '  "validation_criteria": [', "Courier New", 10.5, c_gray)
    add_bullet(json_tf, '    "Endpoint /api/status returns online",', "Courier New", 10.5, c_gray)
    add_bullet(json_tf, '    "File compiles without SyntaxError"', "Courier New", 10.5, c_gray)
    add_bullet(json_tf, '  ],', "Courier New", 10.5, c_gray)
    add_bullet(json_tf, '  "cores_required": ["Core 00", "Core 01", "Core 02"]', "Courier New", 10.5, c_gray)
    add_bullet(json_tf, "}", "Courier New", 10.5, c_gray)
    add_bullet(json_tf, "", "Courier New", 10.5, c_gray)
    add_bullet(json_tf, "INTENT COMPILATION: SUCCESSFUL (200 OK)", "Courier New", 10, c_lime, True)

    add_slide_footer(s6, 6)

    # ----------------------------------------------------
    # SLIDE 7: MODE 02: DEVCORE MODE
    # ----------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7)
    add_slide_header(s7, "Mode 02: DevCore Mode (Live Telemetry)", "SYSTEM MODES OVERVIEW")
    
    # Left Column
    draw_card(s7, Inches(0.6), Inches(1.5), Inches(5.8), Inches(5.2), c_card, c_border)
    tx = add_textbox(s7, Inches(0.9), Inches(1.7), Inches(5.2), Inches(4.8))
    tf = tx.text_frame
    tf.paragraphs[0].text = "Real-Time Telemetry & Monitors"
    tf.paragraphs[0].font.name = "Consolas"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = c_white
    tf.paragraphs[0].font.bold = True
    add_bullet(tf, "Live Dashboards: Exposes internal logs and status codes across the entire co-processor grid.", "Calibri", 12, c_gray)
    add_bullet(tf, "Activity Monitors: Shows color-coded states for each core (IDLE in dark gray, ACTIVE in pulsing lime, ERROR in blinking red).", "Calibri", 12, c_gray)
    add_bullet(tf, "Low Latency: Updates status and file writing actions utilizing short-polling WebSocket telemetry connections.", "Calibri", 12, c_gray)

    # Right Column - Console Interface
    draw_card(s7, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.2), RGBColor(12, 14, 18), c_border)
    add_textbox(s7, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.4), "// CORE PROCESSOR REGISTERS (DEVCORE_MODE)", "Courier New", 11, c_cyan, True)
    
    c_box = add_textbox(s7, Inches(7.1), Inches(2.2), Inches(5.3), Inches(4.2))
    c_tf = c_box.text_frame
    c_tf.paragraphs[0].text = "Planner   [Core 00]  |  STATE: IDLE        |  REG: 0x0"
    c_tf.paragraphs[0].font.name = "Courier New"
    c_tf.paragraphs[0].font.size = Pt(10)
    c_tf.paragraphs[0].font.color.rgb = c_gray
    add_bullet(c_tf, "Coder     [Core 01]  |  STATE: WRITING     |  REG: 0x93FF", "Courier New", 10, c_lime, True)
    add_bullet(c_tf, "Tester    [Core 02]  |  STATE: WAITING     |  REG: 0x0", "Courier New", 10, c_gray)
    add_bullet(c_tf, "Corrector [Core 03]  |  STATE: WAITING     |  REG: 0x0", "Courier New", 10, c_gray)
    add_bullet(c_tf, "Deployer  [Core 04]  |  STATE: WAITING     |  REG: 0x0", "Courier New", 10, c_gray)
    add_bullet(c_tf, "Ledger    [Core 05]  |  STATE: SYNCED      |  REG: 0x4A1", "Courier New", 10, c_cyan)
    add_bullet(c_tf, "", "Courier New", 10, c_gray)
    add_bullet(c_tf, "FILE_BUFFER_TARGET: file:///app/logistics_tracker.py", "Courier New", 10, c_gray)
    add_bullet(c_tf, "LAST_WRITE_OP: write_block (lines 12-45 succeeded)", "Courier New", 10, c_lime)

    add_slide_footer(s7, 7)

    # ----------------------------------------------------
    # SLIDE 8: MODE 03: CHAOS MODE (WITH CHAOS/BROKEN CODE IMAGE)
    # ----------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8)
    add_slide_header(s8, "Mode 03: Chaos Mode (Disruption Testing)", "SYSTEM MODES OVERVIEW")
    
    # Left Top Column - Text
    draw_card(s8, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.4), c_card, c_border)
    tx = add_textbox(s8, Inches(0.9), Inches(1.6), Inches(5.2), Inches(2.2))
    tf = tx.text_frame
    tf.paragraphs[0].text = "Resilience and Stress-Testing"
    tf.paragraphs[0].font.name = "Consolas"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = c_white
    tf.paragraphs[0].font.bold = True
    add_bullet(tf, "Disruption Simulator: Injects deliberate logic faults, file write locks, and database connection timeouts.", "Calibri", 11.5, c_gray)
    add_bullet(tf, "Healing Assessment: Tests if Core 03 (Corrector) can intercept raw compiler stack traces and fix errors autonomously.", "Calibri", 11.5, c_gray)

    # Left Bottom Column - Image
    draw_card(s8, Inches(0.6), Inches(4.1), Inches(5.8), Inches(2.6), c_bg, c_border)
    try_add_picture(s8, img_broken, Inches(0.65), Inches(4.15), Inches(5.7), Inches(2.5))

    # Right Column - Chaos Active Console
    draw_card(s8, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.2), RGBColor(22, 12, 12), c_red)
    add_textbox(s8, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.4), "!! WARNING: CHAOS MONKEY STRESS TRIGGERED !!", "Courier New", 11, c_red, True)
    
    ch_box = add_textbox(s8, Inches(7.1), Inches(2.2), Inches(5.3), Inches(4.2))
    ch_tf = ch_box.text_frame
    ch_tf.paragraphs[0].text = "[CHAOS_SIM] State changed to ACTIVE."
    ch_tf.paragraphs[0].font.name = "Courier New"
    ch_tf.paragraphs[0].font.size = Pt(10)
    ch_tf.paragraphs[0].font.color.rgb = c_red
    add_bullet(ch_tf, "[CHAOS_SIM] Injecting ZeroDivisionError into dashboard_analytics.py", "Courier New", 10, c_red)
    add_bullet(ch_tf, "[CORE 02] Running test verification...", "Courier New", 10, c_gray)
    add_bullet(ch_tf, "[CORE 02] TEST FAILED: ZeroDivisionError at line 42.", "Courier New", 10, c_red, True)
    add_bullet(ch_tf, "[CORE 03] Intercepting traceback exception...", "Courier New", 10, c_lime)
    add_bullet(ch_tf, "[CORE 03] Hotfix applied: wrapped line 42 in denominator checks.", "Courier New", 10, c_lime, True)
    add_bullet(ch_tf, "[CORE 02] Re-running verification checks...", "Courier New", 10, c_gray)
    add_bullet(ch_tf, "[CORE 02] VERIFICATION COMPLETED (Exit Code 0).", "Courier New", 10, c_lime)

    add_slide_footer(s8, 8)

    # ----------------------------------------------------
    # SLIDE 9: MODE 04: CHRONOS MODE (WITH TIMELINE IMAGE)
    # ----------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9)
    add_slide_header(s9, "Mode 04: Chronos Mode (Ledger Timeline)", "SYSTEM MODES OVERVIEW")
    
    # Left Top Column - Text
    draw_card(s9, Inches(0.6), Inches(1.5), Inches(5.8), Inches(2.4), c_card, c_border)
    tx = add_textbox(s9, Inches(0.9), Inches(1.6), Inches(5.2), Inches(2.2))
    tf = tx.text_frame
    tf.paragraphs[0].text = "Time-Travel State Restoration"
    tf.paragraphs[0].font.name = "Consolas"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = c_white
    tf.paragraphs[0].font.bold = True
    add_bullet(tf, "Workspace Snapshots: Records full checklists, parameters, and physical source file states at every compilation node.", "Calibri", 11.5, c_gray)
    add_bullet(tf, "Rollbacks: Triggers zero-loss workspace restoration to any historical checkpoint hash, instantly resetting processor states.", "Calibri", 11.5, c_gray)

    # Left Bottom Column - Image
    draw_card(s9, Inches(0.6), Inches(4.1), Inches(5.8), Inches(2.6), c_bg, c_border)
    try_add_picture(s9, img_clock, Inches(0.65), Inches(4.15), Inches(5.7), Inches(2.5))

    # Right Column - Timeline
    draw_card(s9, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.2), RGBColor(16, 18, 16), c_lime)
    add_textbox(s9, Inches(7.1), Inches(1.7), Inches(5.3), Inches(0.4), "⏳ CHRONOS TIME-TRAVEL JOURNAL", "Courier New", 11, c_lime, True)
    
    t_box = add_textbox(s9, Inches(7.1), Inches(2.2), Inches(5.3), Inches(4.2))
    t_tf = t_box.text_frame
    t_tf.paragraphs[0].text = "Snapshot #189  |  STATE: PLANNING  |  (11:22:29)"
    t_tf.paragraphs[0].font.name = "Courier New"
    t_tf.paragraphs[0].font.size = Pt(10)
    t_tf.paragraphs[0].font.color.rgb = c_gray
    add_bullet(t_tf, "Snapshot #190  |  STATE: CODING    |  (11:22:37)", "Courier New", 10, c_gray)
    add_bullet(t_tf, "Snapshot #191  |  STATE: TESTED    |  (11:22:40)", "Courier New", 10, c_gray)
    add_bullet(t_tf, "Snapshot #192  |  STATE: DEPLOYED  |  (11:22:42) // LATEST", "Courier New", 10, c_cyan)
    add_bullet(t_tf, "", "Courier New", 10, c_gray)
    add_bullet(t_tf, "[ACTION] ROLLBACK INITIATED TO SNAPSHOT #189", "Courier New", 10, c_red, True)
    add_bullet(t_tf, "[LEDGER] Reverting source files to Snapshot #189 hash...", "Courier New", 10, c_gray)
    add_bullet(t_tf, "[LEDGER] System registers reverted. Core states reset.", "Courier New", 10, c_lime, True)

    add_slide_footer(s9, 9)

    # ----------------------------------------------------
    # SLIDE 10: UNDER THE HOOD (WITH CENTER NETWORK FLOW IMAGE)
    # ----------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10)
    add_slide_header(s10, "Under the Hood: LangGraph & SQLite Ledger", "TECHNICAL IMPLEMENTATION")
    
    # Left Card
    draw_card(s10, Inches(0.6), Inches(1.5), Inches(3.8), Inches(5.2), c_card, c_border)
    add_textbox(s10, Inches(0.8), Inches(1.7), Inches(3.4), Inches(0.4), "01 / LANGGRAPH STATE", "Courier New", 12, c_lime, True)
    tx = add_textbox(s10, Inches(0.8), Inches(2.2), Inches(3.4), Inches(4.2))
    tf = tx.text_frame
    tf.paragraphs[0].text = "Message Routing"
    tf.paragraphs[0].font.name = "Consolas"
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = c_white
    tf.paragraphs[0].font.bold = True
    add_bullet(tf, "Coordinates all tasks using a central dictionary state containing target codes and test logs.", "Calibri", 10.5, c_gray)
    add_bullet(tf, "Cores run as isolated execution nodes in a closed cycle.", "Calibri", 10.5, c_gray)

    # Center Image Card
    draw_card(s10, Inches(4.8), Inches(1.5), Inches(3.73), Inches(5.2), c_bg, c_border)
    try_add_picture(s10, img_core, Inches(4.85), Inches(1.55), Inches(3.63), Inches(5.1))

    # Right Card
    draw_card(s10, Inches(8.93), Inches(1.5), Inches(3.8), Inches(5.2), c_card, c_border)
    add_textbox(s10, Inches(9.13), Inches(1.7), Inches(3.4), Inches(0.4), "02 / SQLite DUAL LEDGER", "Courier New", 12, c_cyan, True)
    tx2 = add_textbox(s10, Inches(9.13), Inches(2.2), Inches(3.4), Inches(4.2))
    tf2 = tx2.text_frame
    tf2.paragraphs[0].text = "Timeline Checking"
    tf2.paragraphs[0].font.name = "Consolas"
    tf2.paragraphs[0].font.size = Pt(13)
    tf2.paragraphs[0].font.color.rgb = c_white
    tf2.paragraphs[0].font.bold = True
    add_bullet(tf2, "Writes commits locally to a time-series journal in `hydradb.db` at every single core check.", "Calibri", 10.5, c_gray)
    add_bullet(tf2, "Enables instant git-like diff comparisons and rollback queries.", "Calibri", 10.5, c_gray)

    add_slide_footer(s10, 10)

    # ----------------------------------------------------
    # SLIDE 11: DEVELOPER & BUSINESS VALUE
    # ----------------------------------------------------
    s11 = prs.slides.add_slide(blank_layout)
    set_bg(s11)
    add_slide_header(s11, "Value Analysis: Reliability & Development Speed", "BUSINESS AND OPERATION IMPACT")
    
    # 3 Cards
    # Card 1
    draw_card(s11, start_x, Inches(1.5), card_w, card_h, c_card, c_border)
    add_textbox(s11, start_x + Inches(0.2), Inches(1.7), card_w - Inches(0.4), Inches(0.4), "99% LESS MANUAL DEBUGGING", "Courier New", 12, c_lime, True)
    txt1 = add_textbox(s11, start_x + Inches(0.2), Inches(2.2), card_w - Inches(0.4), Inches(4.2))
    tf1 = txt1.text_frame
    tf1.paragraphs[0].text = "Autonomous Hotfixing"
    tf1.paragraphs[0].font.name = "Consolas"
    tf1.paragraphs[0].font.size = Pt(13)
    tf1.paragraphs[0].font.color.rgb = c_white
    tf1.paragraphs[0].font.bold = True
    add_bullet(tf1, "Core 03 (Corrector) intercepts traceback exceptions and writes patches immediately.", "Calibri", 11, c_gray)
    add_bullet(tf1, "Reduces development cycles by bypassing continuous manual code adjustments and testing.", "Calibri", 11, c_gray)

    # Card 2
    draw_card(s11, start_x + card_w + spacing, Inches(1.5), card_w, card_h, c_card, c_border)
    add_textbox(s11, start_x + card_w + spacing + Inches(0.2), Inches(1.7), card_w - Inches(0.4), Inches(0.4), "ZERO RISK CODE MUTATIONS", "Courier New", 12, c_cyan, True)
    txt2 = add_textbox(s11, start_x + card_w + spacing + Inches(0.2), Inches(2.2), card_w - Inches(0.4), Inches(4.2))
    tf2 = txt2.text_frame
    tf2.paragraphs[0].text = "Fail-Safe Rollbacks"
    tf2.paragraphs[0].font.name = "Consolas"
    tf2.paragraphs[0].font.size = Pt(13)
    tf2.paragraphs[0].font.color.rgb = c_white
    tf2.paragraphs[0].font.bold = True
    add_bullet(tf2, "Chronos Mode makes destructive software trials risk-free.", "Calibri", 11, c_gray)
    add_bullet(tf2, "Allows developers to safely test extreme prompts, secure in the knowledge that any historical snapshot can be restored instantly.", "Calibri", 11, c_gray)

    # Card 3
    draw_card(s11, start_x + (card_w + spacing)*2, Inches(1.5), card_w, card_h, c_card, c_border)
    add_textbox(s11, start_x + (card_w + spacing)*2 + Inches(0.2), Inches(1.7), card_w - Inches(0.4), Inches(0.4), "TRANSPARENT SYSTEM TRAILS", "Courier New", 12, c_white, True)
    txt3 = add_textbox(s11, start_x + (card_w + spacing)*2 + Inches(0.2), Inches(2.2), card_w - Inches(0.4), Inches(4.2))
    tf3 = txt3.text_frame
    tf3.paragraphs[0].text = "Audit-Ready Logs"
    tf3.paragraphs[0].font.name = "Consolas"
    tf3.paragraphs[0].font.size = Pt(13)
    tf3.paragraphs[0].font.color.rgb = c_white
    tf3.paragraphs[0].font.bold = True
    add_bullet(tf3, "All agent prompts, checklist outputs, and file writes are logged systematically to SQLite.", "Calibri", 11, c_gray)
    add_bullet(tf3, "Provides comprehensive compliance records, making multi-agent operations fully transparent and explainable.", "Calibri", 11, c_gray)

    add_slide_footer(s11, 11)

    # ----------------------------------------------------
    # SLIDE 12: FUTURE ROADMAP & CONCLUSION
    # ----------------------------------------------------
    s12 = prs.slides.add_slide(blank_layout)
    set_bg(s12)
    add_slide_header(s12, "Roadmap: Scaling Agent Runtimes", "FUTURE OUTLOOK AND CONCLUSION")
    
    # Left Column
    draw_card(s12, Inches(0.6), Inches(1.5), Inches(6.8), Inches(5.2), c_card, c_border)
    tx = add_textbox(s12, Inches(0.9), Inches(1.7), Inches(6.2), Inches(4.8))
    tf = tx.text_frame
    tf.paragraphs[0].text = "Future Project Milestones"
    tf.paragraphs[0].font.name = "Consolas"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = c_white
    tf.paragraphs[0].font.bold = True
    add_bullet(tf, "Kratos OS v1.5 [Kubernetes Integration]: Scaled test runners to isolated, transient container pods for parallel compilation verification.", "Calibri", 12, c_gray)
    add_bullet(tf, "Kratos OS v2.0 [Hardware Enclaves]: Compilation and file writes executed inside Intel SGX or AWS Nitro Enclaves to guarantee zero-trust code environments.", "Calibri", 12, c_gray)
    add_bullet(tf, "Kratos OS v2.5 [Consensus Engine]: Multi-agent code peer review. Codes must be approved by three separate LLM verifiers before deployment.", "Calibri", 12, c_gray)

    # Right Column - Summary Box
    draw_card(s12, Inches(7.8), Inches(1.5), Inches(4.9), Inches(5.2), c_card, c_lime)
    add_textbox(s12, Inches(8.1), Inches(1.7), Inches(4.3), Inches(0.4), "KRATOS OS SANDBOX ACTIVE", "Courier New", 12, c_lime, True)
    
    sum_box = add_textbox(s12, Inches(8.1), Inches(2.2), Inches(4.3), Inches(4.2))
    sum_tf = sum_box.text_frame
    sum_tf.paragraphs[0].text = "System Summary:"
    sum_tf.paragraphs[0].font.name = "Consolas"
    sum_tf.paragraphs[0].font.size = Pt(13)
    sum_tf.paragraphs[0].font.color.rgb = c_white
    sum_tf.paragraphs[0].font.bold = True
    add_bullet(sum_tf, "✔ 6-Core Engine is active and stable.", "Calibri", 11, c_gray)
    add_bullet(sum_tf, "✔ Prompt, DevCore, Chaos, and Chronos Modes fully operational.", "Calibri", 11, c_gray)
    add_bullet(sum_tf, "✔ Secure SQLite timeline persistence active.", "Calibri", 11, c_gray)
    add_bullet(sum_tf, "✔ Ready for natural language intent scaffolding.", "Calibri", 11, c_lime, True)

    add_slide_footer(s12, 12)

    # Save to disk
    prs.save("KratosOS_Briefing_Deck_v2.pptx")
    print("PowerPoint presentation generated successfully: KratosOS_Briefing_Deck_v2.pptx")

if __name__ == "__main__":
    create_deck()
